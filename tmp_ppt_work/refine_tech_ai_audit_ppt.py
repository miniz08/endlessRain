# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SRC = Path(r"D:\$RRBBDQK\tmp_ppt_work\defense_9page_source.pptx")
OUT = Path(r"D:\$RRBBDQK\tmp_ppt_work\defense_ppt_tech_ai_audit_refined.pptx")

FONT = "Microsoft YaHei"
TITLE = RGBColor(0x24, 0x38, 0x46)
ACCENT = RGBColor(0x15, 0x71, 0x6E)
MUTED = RGBColor(0x60, 0x6E, 0x76)
LIGHT = RGBColor(0xE4, 0xF7, 0xF3)
LINE = RGBColor(0xB3, 0xD9, 0xD3)


def remove_shape(shape):
    element = shape._element
    element.getparent().remove(element)


def clear_body(slide):
    for shape in list(slide.shapes)[4:]:
        remove_shape(shape)


def set_title(slide, title):
    shape = slide.shapes[2]
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.size = Pt(21)
    p.font.bold = True
    p.font.color.rgb = TITLE


def textbox(slide, x, y, w, h, text="", size=13, color=TITLE, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return shape


def hline(slide, x, y, w, color=LINE, width=1.0):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x + w), Inches(y))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def soft_band(slide, x, y, w, h, text, size=13.5, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = LINE
    shape.line.width = Pt(0.7)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    return shape


def section(slide, x, y, w, title, lines, title_size=15.5, body_size=10.9):
    textbox(slide, x, y, w, 0.3, title, title_size, ACCENT, True)
    hline(slide, x, y + 0.38, w, ACCENT, 1.2)
    yy = y + 0.55
    for line in lines:
        textbox(slide, x, yy, w, 0.25, line, body_size, MUTED)
        yy += 0.31


def number_step(slide, x, y, no, title, desc, desc_w=7.5):
    textbox(slide, x, y, 0.4, 0.27, no, 12.5, ACCENT, True)
    textbox(slide, x + 0.46, y, 1.45, 0.27, title, 12.8, TITLE, True)
    textbox(slide, x + 1.75, y, desc_w, 0.27, desc, 10.6, MUTED)


def metric(slide, x, y, value, label):
    textbox(slide, x, y, 2.0, 0.34, value, 21, ACCENT, True, PP_ALIGN.CENTER)
    textbox(slide, x, y + 0.4, 2.0, 0.24, label, 10.6, MUTED, False, PP_ALIGN.CENTER)


def build_tech_architecture(slide):
    set_title(slide, "总体架构：技术栈与运行方式")
    clear_body(slide)

    textbox(slide, 0.95, 1.02, 10.9, 0.34, "系统基于 Node.js 技术栈实现，通过 Docker 容器运行，对外由 Nginx / 网关统一承接请求。", 15.8, ACCENT, True, PP_ALIGN.CENTER)

    section(
        slide,
        0.95,
        1.72,
        3.15,
        "前端与入口",
        [
            "Nuxt3 / Vue：页面与交互",
            "Nginx：前端容器静态资源与代理",
            "API Gateway：统一 API 入口",
            "WebSocket：聊天连接代理",
        ],
    )
    section(
        slide,
        4.55,
        1.72,
        3.15,
        "后端服务",
        [
            "Node.js + Express：服务接口",
            "TypeScript：主要业务代码",
            "http-proxy-middleware：网关转发",
            "express-rate-limit：入口限流",
        ],
    )
    section(
        slide,
        8.05,
        1.72,
        3.45,
        "数据与部署",
        [
            "MySQL 5.7.44：longseason 数据库",
            "Prisma：数据访问与表结构映射",
            "Docker / Compose：服务镜像与编排",
            "dotenv：环境变量与模型配置",
        ],
    )

    hline(slide, 0.95, 4.35, 10.65)
    textbox(slide, 1.0, 4.78, 1.65, 0.28, "请求路径", 14.8, ACCENT, True)
    textbox(slide, 2.55, 4.78, 8.9, 0.28, "浏览器 -> Nuxt3/Nginx -> api_gateway -> user/blog/ai/chat service -> MySQL", 11.6, MUTED)
    textbox(slide, 1.0, 5.35, 1.65, 0.28, "追踪方式", 14.8, ACCENT, True)
    textbox(slide, 2.55, 5.35, 8.9, 0.28, "网关生成并透传 requestId，业务服务写入推荐记录或审计日志。", 11.6, MUTED)
    soft_band(slide, 1.0, 5.94, 10.65, 0.38, "讲述重点：技术栈服务于三件事：统一入口、服务隔离、请求可追踪。", 13.2)


def build_ai_governance(slide):
    set_title(slide, "AI 治理：外部模型如何接入")
    clear_body(slide)

    textbox(slide, 0.95, 1.02, 10.9, 0.34, "AI 服务不直接绑定某个模型，而是用 OpenAI-compatible Provider 封装外部大模型调用。", 15.8, ACCENT, True, PP_ALIGN.CENTER)

    steps = [
        ("01", "触发", "文章创建后，blog_service 携带内部 token 调用 ai_service；文章先保持待审核。"),
        ("02", "构造 message", "system 规定审核器角色和 JSON 输出；user 放入任务、schema、标签库和文章内容。"),
        ("03", "兼容请求", "按 OpenAI-compatible 风格 POST /chat/completions，body 包含 model、messages、temperature、response_format。"),
        ("04", "轮询 Provider", "从环境变量读取多个国内模型商配置，失败时切换下一个目标或模型。"),
        ("05", "解析规范化", "解析模型 JSON，规范化四项评分，校验标签来自预设标签库。"),
        ("06", "落库使用", "写入 AI 分析表与标签关系表，并更新文章状态、风险等级和审核建议。"),
    ]
    y = 1.85
    for no, title, desc in steps:
        number_step(slide, 1.0, y, no, title, desc, 8.8)
        hline(slide, 1.0, y + 0.42, 10.55)
        y += 0.58

    soft_band(
        slide,
        1.0,
        5.55,
        5.15,
        0.58,
        "message 语义：不是自由聊天，而是“内容审核任务 + 输出格式约束”。",
        12.4,
    )
    soft_band(
        slide,
        6.55,
        5.55,
        5.15,
        0.58,
        "失败策略：外部模型异常时保守降级，文章进入 REVIEW_REQUIRED。",
        12.4,
    )


def build_audit_flow(slide):
    set_title(slide, "安全审计：从方案到追踪")
    clear_body(slide)

    textbox(slide, 0.95, 1.02, 10.9, 0.34, "审计不是单独展示日志，而是设计一条能定位请求、记录行为、支持查询的证据链。", 15.8, ACCENT, True, PP_ALIGN.CENTER)

    textbox(slide, 1.0, 1.78, 2.05, 0.3, "审计方案", 16, ACCENT, True)
    hline(slide, 1.0, 2.19, 2.8, ACCENT, 1.2)
    textbox(slide, 1.0, 2.44, 3.15, 0.28, "网关生成 requestId 并透传", 11.4, MUTED)
    textbox(slide, 1.0, 2.80, 3.15, 0.28, "业务服务做角色与资源归属校验", 11.4, MUTED)
    textbox(slide, 1.0, 3.16, 3.15, 0.28, "关键动作统一写 audit_log", 11.4, MUTED)
    textbox(slide, 1.0, 3.52, 3.15, 0.28, "管理员通过运维页查询", 11.4, MUTED)

    textbox(slide, 4.55, 1.78, 2.05, 0.3, "执行写入", 16, ACCENT, True)
    hline(slide, 4.55, 2.19, 2.8, ACCENT, 1.2)
    textbox(slide, 4.55, 2.44, 3.25, 0.28, "登录、退出、注册等身份动作", 11.4, MUTED)
    textbox(slide, 4.55, 2.80, 3.25, 0.28, "发文、删除、评论、关注等业务动作", 11.4, MUTED)
    textbox(slide, 4.55, 3.16, 3.25, 0.28, "AI 分析与管理员查询动作", 11.4, MUTED)
    textbox(slide, 4.55, 3.52, 3.25, 0.28, "聊天与异常请求记录状态码", 11.4, MUTED)

    textbox(slide, 8.15, 1.78, 2.05, 0.3, "追踪结果", 16, ACCENT, True)
    hline(slide, 8.15, 2.19, 2.8, ACCENT, 1.2)
    textbox(slide, 8.15, 2.44, 3.25, 0.28, "按 requestId 定位一次请求链路", 11.4, MUTED)
    textbox(slide, 8.15, 2.80, 3.25, 0.28, "按用户、路径、动作和状态码筛查", 11.4, MUTED)
    textbox(slide, 8.15, 3.16, 3.25, 0.28, "运维页展示审计列表与摘要", 11.4, MUTED)
    textbox(slide, 8.15, 3.52, 3.25, 0.28, "异常时可回到具体服务排查", 11.4, MUTED)

    hline(slide, 1.0, 4.32, 10.55)
    metric(slide, 1.05, 4.78, "requestId", "请求追踪")
    metric(slide, 3.65, 4.78, "audit_log", "行为证据")
    metric(slide, 6.25, 4.78, "运维页", "查询入口")
    metric(slide, 8.85, 4.78, "状态码", "结果定位")
    soft_band(slide, 1.0, 5.84, 10.65, 0.42, "讲述重点：审计闭环 = 入口标识 -> 业务写入 -> 管理查询 -> 异常定位。", 13.6)


def main():
    prs = Presentation(SRC)
    build_tech_architecture(prs.slides[2])
    build_ai_governance(prs.slides[4])
    build_audit_flow(prs.slides[6])
    prs.save(OUT)
    print(str(OUT))


if __name__ == "__main__":
    main()
