# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SRC = Path(r"D:\$RRBBDQK\tmp_ppt_work\defense_ppt_rebuilt_source.pptx")
OUT = Path(r"D:\$RRBBDQK\tmp_ppt_work\defense_ppt_streamlined.pptx")

FONT = "Microsoft YaHei"
TITLE = RGBColor(0x24, 0x38, 0x46)
ACCENT = RGBColor(0x15, 0x71, 0x6E)
MUTED = RGBColor(0x60, 0x6E, 0x76)
LIGHT = RGBColor(0xE4, 0xF7, 0xF3)
SOFT = RGBColor(0xF8, 0xFC, 0xFB)
LINE = RGBColor(0xB3, 0xD9, 0xD3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def remove_shape(shape):
    element = shape._element
    element.getparent().remove(element)


def clear_body(slide):
    for shape in list(slide.shapes)[4:]:
        remove_shape(shape)


def set_title(slide, title):
    shape = slide.shapes[2]
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.size = Pt(21)
    p.font.bold = True
    p.font.color.rgb = TITLE


def textbox(slide, x, y, w, h, text="", size=14, color=TITLE, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return shape


def paragraph_box(slide, x, y, w, h, title, lines, title_size=16, body_size=12.5):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    for line in lines:
        q = tf.add_paragraph()
        q.text = line
        q.space_before = Pt(4)
        q.font.name = FONT
        q.font.size = Pt(body_size)
        q.font.color.rgb = MUTED
    return shape


def hline(slide, x, y, w, color=LINE, width=1.2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x + w), Inches(y))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def vline(slide, x, y, h, color=LINE, width=1.2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x), Inches(y + h))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def arrow(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = ACCENT
    line.line.width = Pt(1.5)
    try:
        line.line.end_arrowhead = True
    except Exception:
        pass
    return line


def soft_band(slide, x, y, w, h, text, size=17, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = LINE
    shape.line.width = Pt(0.7)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.16)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    return shape


def metric_text(slide, x, y, value, label, note=None):
    textbox(slide, x, y, 2.25, 0.38, value, 24, ACCENT, True, PP_ALIGN.CENTER)
    textbox(slide, x, y + 0.42, 2.25, 0.28, label, 11.5, MUTED, False, PP_ALIGN.CENTER)
    if note:
        textbox(slide, x, y + 0.73, 2.25, 0.25, note, 10.5, MUTED, False, PP_ALIGN.CENTER)


def numbered_line(slide, x, y, no, title, desc):
    textbox(slide, x, y, 0.45, 0.3, no, 14, ACCENT, True, PP_ALIGN.LEFT)
    textbox(slide, x + 0.45, y, 2.0, 0.3, title, 14, TITLE, True, PP_ALIGN.LEFT)
    textbox(slide, x + 2.12, y, 5.6, 0.32, desc, 12.2, MUTED, False, PP_ALIGN.LEFT)


def build_slide_3(slide):
    set_title(slide, "总体架构：统一入口与分层协作")
    clear_body(slide)

    soft_band(slide, 0.78, 1.05, 11.55, 0.55, "对外只有一个访问入口；对内按职责拆分服务，所有关键请求带 requestId。", 17)

    y = 1.96
    sections = [
        ("客户端与前端", "浏览器访问 Nuxt3 前端；前端负责页面展示、交互和调用 API。"),
        ("网关中间层", "api_gateway 统一转发到后端服务，并处理限流、requestId 透传和 WebSocket 代理。"),
        ("业务服务层", "user_service 管身份；blog_service 管内容、互动、推荐、通知；ai_service 管内容分析；chat_service 管私信。"),
        ("数据层", "MySQL longseason 保存业务表、推荐事件、AI 分析结果和审计日志。"),
    ]
    for title, desc in sections:
        textbox(slide, 1.0, y, 2.05, 0.32, title, 15.5, ACCENT, True)
        textbox(slide, 3.0, y, 8.8, 0.35, desc, 13.2, MUTED)
        hline(slide, 1.0, y + 0.46, 10.9)
        y += 0.86

    textbox(slide, 1.0, 5.72, 10.9, 0.45, "讲述重点：架构图不是为了展示服务数量，而是说明请求路径、服务边界和追踪方式。", 14, TITLE, False, PP_ALIGN.CENTER)


def build_slide_4(slide):
    set_title(slide, "设计主线：从发文到追踪")
    clear_body(slide)

    textbox(slide, 0.9, 1.08, 11.1, 0.45, "核心不是把功能摆满，而是让内容从发布、治理、分发到追踪形成闭环。", 18, ACCENT, True, PP_ALIGN.CENTER)
    hline(slide, 1.0, 1.72, 10.9, ACCENT, 1.4)

    steps = [
        ("请求进入", "网关转发、限流，并透传 requestId"),
        ("内容入库", "文章先以待审核状态保存，互动写入业务表"),
        ("AI 治理", "生成四项评分、标签、风险等级和处理决策"),
        ("展示推荐", "先判断能不能展示，再在候选池里排序"),
        ("回写追踪", "用户行为写入推荐事件，关键操作写入审计日志"),
    ]
    x = 1.0
    y = 2.45
    for i, (title, desc) in enumerate(steps, 1):
        textbox(slide, x, y, 0.55, 0.36, f"{i:02d}", 15.5, ACCENT, True)
        textbox(slide, x + 0.62, y, 1.65, 0.34, title, 15.5, TITLE, True)
        textbox(slide, x + 2.1, y, 8.5, 0.34, desc, 12.8, MUTED)
        if i < len(steps):
            hline(slide, x + 0.22, y + 0.48, 10.3)
        y += 0.74

    soft_band(slide, 1.0, 5.92, 10.9, 0.43, "一句话：先管住内容，再分发内容，最后留下证据。", 15.5)


def build_slide_5(slide):
    set_title(slide, "服务划分：三层职责")
    clear_body(slide)

    textbox(slide, 0.95, 1.1, 10.9, 0.42, "六个服务不是平铺罗列，而是按访问路径分成前端、网关中间层和后端服务。", 17, ACCENT, True, PP_ALIGN.CENTER)

    headers = [
        ("前端服务", 0.95, 2.0, 2.55),
        ("网关中间服务", 4.05, 2.0, 3.0),
        ("后端服务", 7.55, 2.0, 3.75),
    ]
    for title, x, y, w in headers:
        textbox(slide, x, y, w, 0.32, title, 16.5, ACCENT, True)
        hline(slide, x, y + 0.42, w, ACCENT, 1.3)

    paragraph_box(slide, 0.95, 2.62, 2.8, 1.4, "mofukaze", ["Nuxt3 前端", "页面展示、发文、通知、私信与运维入口"], 15.5, 11.3)

    paragraph_box(slide, 4.05, 2.62, 3.0, 1.4, "api_gateway", ["统一 API 入口", "路由转发、限流、requestId 透传", "WebSocket 代理与审计查询入口"], 15.5, 11.1)

    paragraph_box(slide, 7.55, 2.62, 1.85, 1.42, "user_service", ["注册登录", "会话与角色"], 14.2, 10.7)
    paragraph_box(slide, 9.65, 2.62, 2.15, 1.42, "blog_service", ["内容互动", "推荐与通知"], 14.2, 10.7)
    paragraph_box(slide, 7.55, 4.33, 1.85, 1.42, "ai_service", ["评分标签", "风险与 Provider"], 14.2, 10.7)
    paragraph_box(slide, 9.65, 4.33, 2.15, 1.42, "chat_service", ["私信线程", "消息与 WebSocket"], 14.2, 10.7)

    vline(slide, 3.8, 1.95, 3.95)
    vline(slide, 7.3, 1.95, 3.95)
    soft_band(slide, 1.0, 6.05, 10.9, 0.4, "讲述重点：前端负责呈现，网关负责入口，后端服务按业务变化点拆分。", 14.5)


def build_slide_6(slide):
    set_title(slide, "AI 治理：结构化结果落库")
    clear_body(slide)

    textbox(slide, 0.95, 1.1, 10.9, 0.42, "AI 调用不是为了得到一段回复，而是为了生成后续业务可使用的结构化字段。", 17, ACCENT, True, PP_ALIGN.CENTER)

    numbered_line(slide, 1.0, 2.0, "01", "触发", "文章创建后，blog_service 调用 ai_service，文章先保持待审核。")
    numbered_line(slide, 1.0, 2.72, "02", "构造", "真实 Provider 使用 messages 风格：system 约束 + user 任务 + 输出 schema + 标签库。")
    numbered_line(slide, 1.0, 3.44, "03", "调用", "Provider 层屏蔽 Mock 和真实模型差异，真实模型可按配置轮询。")
    numbered_line(slide, 1.0, 4.16, "04", "规范化", "解析 JSON，限制分数范围，校验标签必须来自预设标签库。")
    numbered_line(slide, 1.0, 4.88, "05", "落库", "写入 article_ai_analysis、article_ai_tag_on_article，并更新文章状态。")

    soft_band(slide, 1.0, 5.82, 5.25, 0.48, "失败策略：模型异常不放行，文章进入 REVIEW_REQUIRED。", 13.8)
    soft_band(slide, 6.6, 5.82, 5.25, 0.48, "业务影响：展示控制、推荐排序、作者评分与通知。", 13.8)


def build_slide_7(slide):
    set_title(slide, "推荐机制：公式与流程分开")
    clear_body(slide)

    textbox(slide, 0.95, 1.05, 10.9, 0.38, "先过滤可展示候选，再用可解释公式排序；推荐流读取已落库 AI 结果。", 17, ACCENT, True, PP_ALIGN.CENTER)

    soft_band(
        slide,
        0.95,
        1.75,
        5.15,
        1.95,
        "推荐总分 =\n"
        "tagMatch * 4 + authorAffinity * 2\n"
        "+ authorQuality * 1.2 + contentQuality + freshness\n"
        "- riskPenalty - seenPenalty",
        15.6,
    )
    textbox(slide, 0.95, 3.95, 5.15, 0.36, "公式含义：兴趣匹配权重最高，但不会绕过质量、风险和重复曝光控制。", 12.6, MUTED, False, PP_ALIGN.CENTER)

    textbox(slide, 6.65, 1.75, 4.8, 0.35, "推荐流程", 17, ACCENT, True)
    hline(slide, 6.65, 2.18, 4.8, ACCENT, 1.3)
    flow = [
        ("候选过滤", "PUBLISHED / LOW_PRIORITY，且 legalityScore >= 40"),
        ("读取信号", "用户画像、AI 标签、作者质量、14 天统计、已看记录"),
        ("计算排序", "按公式得到 total score，分数相同按发布时间"),
        ("返回并记录", "返回推荐列表，写 reco_request_log、曝光和 seen"),
    ]
    y = 2.5
    for i, (title, desc) in enumerate(flow, 1):
        textbox(slide, 6.65, y, 0.38, 0.28, str(i), 13, ACCENT, True)
        textbox(slide, 7.05, y, 1.35, 0.28, title, 13, TITLE, True)
        textbox(slide, 8.35, y, 3.4, 0.28, desc, 10.8, MUTED)
        y += 0.72

    soft_band(slide, 1.0, 5.83, 10.9, 0.45, "关键边界：推荐请求日志是 reco_request_log，不等同于审计日志 audit_log。", 14.2)


def build_slide_8(slide):
    set_title(slide, "画像反馈：行为改变下一次推荐")
    clear_body(slide)

    textbox(slide, 0.95, 1.05, 10.9, 0.42, "画像不是用户自己填写，而是由推荐事件持续计算出来。", 18, ACCENT, True, PP_ALIGN.CENTER)

    y = 2.05
    flow = [
        ("推荐展示", "返回列表时自动写入曝光事件，并更新已看记录。"),
        ("用户行为", "点击、停留、阅读完成、点赞、评论、隐藏、举报等形成反馈。"),
        ("事件入库", "统一写入 reco_event，保留场景、位置、用户、文章和 requestId。"),
        ("画像刷新", "最近 90 天、最多 500 条事件参与计算，旧行为按时间衰减。"),
        ("影响排序", "形成 tagVector 与 authorAffinity，参与下一次推荐打分。"),
    ]
    for i, (title, desc) in enumerate(flow, 1):
        textbox(slide, 1.0, y, 0.52, 0.3, f"{i}", 15, ACCENT, True)
        textbox(slide, 1.58, y, 1.55, 0.3, title, 15, TITLE, True)
        textbox(slide, 3.0, y, 8.5, 0.32, desc, 12.4, MUTED)
        if i < len(flow):
            hline(slide, 1.0, y + 0.47, 10.9)
        y += 0.74

    textbox(slide, 1.15, 5.88, 5.0, 0.32, "正向行为提高相关标签和作者权重", 13.5, ACCENT, True, PP_ALIGN.CENTER)
    textbox(slide, 6.6, 5.88, 5.0, 0.32, "隐藏、举报降低相似内容权重", 13.5, ACCENT, True, PP_ALIGN.CENTER)


def build_slide_9(slide):
    set_title(slide, "安全审计：从请求到证据")
    clear_body(slide)

    textbox(slide, 0.95, 1.05, 10.9, 0.42, "审计不是“日志很多”，而是关键行为能按请求链路追踪。", 18, ACCENT, True, PP_ALIGN.CENTER)

    paragraph_box(slide, 1.0, 1.95, 3.3, 1.35, "入口控制", ["网关限流与 requestId 透传", "WebSocket 请求同样经过网关"], 16, 11.8)
    paragraph_box(slide, 4.75, 1.95, 3.3, 1.35, "业务校验", ["角色判断", "文章、评论、聊天等资源归属校验"], 16, 11.8)
    paragraph_box(slide, 8.5, 1.95, 3.2, 1.35, "证据保留", ["audit_log 记录主体、路径、动作、结果", "管理员运维页查询"], 16, 11.8)
    hline(slide, 1.0, 3.62, 10.7)

    textbox(slide, 1.0, 4.1, 2.2, 0.3, "记录范围", 15.5, ACCENT, True)
    textbox(slide, 3.0, 4.1, 8.8, 0.3, "登录、发文、删除、AI 分析、评论、关注、聊天、异常请求和管理员查询。", 12.6, MUTED)
    textbox(slide, 1.0, 4.78, 2.2, 0.3, "闭环方式", 15.5, ACCENT, True)
    textbox(slide, 3.0, 4.78, 8.8, 0.3, "requestId 定位请求，audit_log 保留行为证据，运维页提供查询入口。", 12.6, MUTED)

    soft_band(slide, 1.0, 5.85, 10.9, 0.42, "一句话：网关负责入口，业务服务负责边界，数据库保留证据。", 14.8)


def build_slide_10(slide):
    set_title(slide, "测试证据：闭环是否跑通")
    clear_body(slide)

    textbox(slide, 0.95, 1.05, 10.9, 0.42, "测试不只证明“能跑”，还证明问题被发现、修复并复验。", 18, ACCENT, True, PP_ALIGN.CENTER)

    metric_text(slide, 0.95, 1.9, "84 / 82", "综合黑盒测试", "通过率 97.62%")
    metric_text(slide, 3.75, 1.9, "7 / 7", "定向复验", "缺陷闭环")
    metric_text(slide, 6.55, 1.9, "13 / 13", "推荐专项", "画像差异有效")
    metric_text(slide, 9.35, 1.9, "160 / 16", "推荐并发复验", "全 200")
    hline(slide, 0.95, 3.18, 10.7)

    paragraph_box(slide, 1.0, 3.58, 3.4, 1.15, "发现的问题", ["注册角色提权风险", "WebSocket 网关转发", "推荐并发写入冲突"], 15.8, 11.4)
    paragraph_box(slide, 4.85, 3.58, 3.4, 1.15, "推荐专项数据", ["3 类画像用户", "36 作者 / 291 文章", "1021 条推荐事件"], 15.8, 11.4)
    paragraph_box(slide, 8.7, 3.58, 3.1, 1.15, "结论边界", ["验证原型闭环", "不夸大为生产容量"], 15.8, 11.4)

    soft_band(slide, 1.0, 5.82, 10.9, 0.45, "Mock Provider 保证系统测试可重复；真实模型接入能力另做补充验证。", 14.5)


def build_slide_11(slide):
    set_title(slide, "总结与致谢")
    clear_body(slide)

    textbox(slide, 1.0, 1.12, 10.9, 0.45, "完成的是一个可运行、可追踪、可解释的社交平台原型。", 19, ACCENT, True, PP_ALIGN.CENTER)
    hline(slide, 1.0, 1.78, 10.9, ACCENT, 1.4)

    items = [
        ("微服务架构", "前端、网关、用户、内容、AI、聊天按职责拆分。"),
        ("内容治理", "文章发布后通过 AI 评分、标签和风险决策控制展示状态。"),
        ("推荐闭环", "画像、推荐事件、排序公式和回写记录形成可解释推荐流程。"),
        ("审计闭环", "requestId、audit_log 和运维页支撑关键行为追踪。"),
    ]
    y = 2.35
    for title, desc in items:
        textbox(slide, 1.1, y, 2.0, 0.32, title, 15.5, ACCENT, True)
        textbox(slide, 3.05, y, 8.7, 0.32, desc, 12.8, MUTED)
        y += 0.62

    textbox(slide, 1.1, 5.18, 2.0, 0.32, "后续优化", 15.5, ACCENT, True)
    textbox(slide, 3.05, 5.18, 8.7, 0.32, "性能优化、单元测试、人工复核和长期推荐效果评估。", 12.8, MUTED)
    textbox(slide, 1.35, 5.83, 10.4, 0.55, "谢谢各位老师", 28, ACCENT, True, PP_ALIGN.CENTER)


def main():
    prs = Presentation(SRC)
    builders = {
        3: build_slide_3,
        4: build_slide_4,
        5: build_slide_5,
        6: build_slide_6,
        7: build_slide_7,
        8: build_slide_8,
        9: build_slide_9,
        10: build_slide_10,
        11: build_slide_11,
    }
    for slide_no, builder in builders.items():
        builder(prs.slides[slide_no - 1])
    prs.save(OUT)
    print(str(OUT))


if __name__ == "__main__":
    main()
