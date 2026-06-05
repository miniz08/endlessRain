# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SRC = Path(r"D:\$RRBBDQK\tmp_ppt_work\current_stream_source.pptx")
OUT = Path(r"D:\$RRBBDQK\tmp_ppt_work\defense_ppt_merged_reco_test.pptx")

FONT = "Microsoft YaHei"
TITLE = RGBColor(0x24, 0x38, 0x46)
ACCENT = RGBColor(0x15, 0x71, 0x6E)
MUTED = RGBColor(0x60, 0x6E, 0x76)
LIGHT = RGBColor(0xE4, 0xF7, 0xF3)
LINE = RGBColor(0xB3, 0xD9, 0xD3)


def remove_shape(shape):
    element = shape._element
    element.getparent().remove(element)


def clear_body(slide):
    for shape in list(slide.shapes)[4:]:
        remove_shape(shape)


def set_title(slide, title):
    shape = slide.shapes[2]
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.size = Pt(21)
    p.font.bold = True
    p.font.color.rgb = TITLE


def textbox(slide, x, y, w, h, text="", size=13, color=TITLE, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return shape


def section(slide, x, y, w, title, lines, title_size=15.4, body_size=10.9):
    textbox(slide, x, y, w, 0.3, title, title_size, ACCENT, True)
    hline(slide, x, y + 0.38, w, ACCENT, 1.2)
    yy = y + 0.55
    for line in lines:
        textbox(slide, x, yy, w, 0.25, line, body_size, MUTED)
        yy += 0.31


def hline(slide, x, y, w, color=LINE, width=1.0):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x + w), Inches(y))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def soft_band(slide, x, y, w, h, text, size=14.5, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = LINE
    shape.line.width = Pt(0.7)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    return shape


def number_step(slide, x, y, no, title, desc):
    textbox(slide, x, y, 0.34, 0.27, no, 12.5, ACCENT, True)
    textbox(slide, x + 0.38, y, 1.25, 0.27, title, 12.5, TITLE, True)
    textbox(slide, x + 1.45, y, 3.75, 0.27, desc, 10.3, MUTED)


def metric(slide, x, y, value, label):
    textbox(slide, x, y, 2.05, 0.34, value, 21, ACCENT, True, PP_ALIGN.CENTER)
    textbox(slide, x, y + 0.4, 2.05, 0.24, label, 10.7, MUTED, False, PP_ALIGN.CENTER)


def delete_slide(prs, index_zero_based):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    r_id = slides[index_zero_based].rId
    prs.part.drop_rel(r_id)
    xml_slides.remove(slides[index_zero_based])


def build_merged_recommendation(slide):
    set_title(slide, "推荐机制：AI 分析与用户画像共同作用")
    clear_body(slide)

    textbox(slide, 0.95, 1.05, 10.9, 0.36, "推荐不是单独的公式，而是内容侧信号与兴趣侧信号共同进入排序。", 16.5, ACCENT, True, PP_ALIGN.CENTER)

    section(
        slide,
        0.95,
        1.72,
        3.35,
        "AI 分析提供内容侧信号",
        [
            "AI 标签：用于和 tagVector 匹配",
            "四项评分：参与内容质量分",
            "合法性分：参与风险过滤与惩罚",
            "文章状态：决定能否进入候选池",
        ],
    )
    section(
        slide,
        8.25,
        1.72,
        3.35,
        "用户画像提供兴趣侧信号",
        [
            "reco_event 记录用户行为",
            "最近 90 天、最多 500 条事件",
            "生成 tagVector 与 authorAffinity",
            "旧行为按时间衰减",
        ],
    )

    soft_band(
        slide,
        4.58,
        1.92,
        3.25,
        1.55,
        "推荐总分\n"
        "tagMatch * 4 + authorAffinity * 2\n"
        "+ authorQuality * 1.2 + contentQuality + freshness\n"
        "- riskPenalty - seenPenalty",
        12.1,
    )
    textbox(slide, 4.52, 3.63, 3.45, 0.25, "先过滤可展示候选，再在候选池中排序。", 10.8, MUTED, False, PP_ALIGN.CENTER)

    hline(slide, 0.95, 4.18, 10.75)
    textbox(slide, 0.95, 4.45, 2.0, 0.3, "推荐闭环流程", 15.2, ACCENT, True)
    steps = [
        ("1", "候选过滤", "PUBLISHED / LOW_PRIORITY 且 legalityScore >= 40"),
        ("2", "读取信号", "AI 标签与评分 + 用户画像 + 14 天统计 + 已看记录"),
        ("3", "计算排序", "公式得到 total score，输出推荐列表"),
        ("4", "行为回写", "曝光、点击、停留、点赞、评论、隐藏、举报写入 reco_event"),
        ("5", "画像刷新", "更新 tagVector / authorAffinity，影响下一次推荐"),
    ]
    y = 4.9
    for no, title, desc in steps:
        number_step(slide, 1.05, y, no, title, desc)
        y += 0.31

    soft_band(slide, 6.7, 5.18, 4.9, 0.58, "关键表达：AI 决定内容质量与风险，画像决定用户兴趣，两者共同决定排序。", 12.6)


def build_testing_flow(slide):
    set_title(slide, "系统测试：从方案到结果")
    clear_body(slide)

    textbox(slide, 0.95, 1.05, 10.9, 0.34, "测试页按执行过程讲：先说明怎么测，再说明发现了什么，最后说明复验结果。", 16.2, ACCENT, True, PP_ALIGN.CENTER)

    textbox(slide, 1.0, 1.82, 2.05, 0.3, "测试方案", 16, ACCENT, True)
    hline(slide, 1.0, 2.23, 2.7, ACCENT, 1.2)
    textbox(slide, 1.0, 2.48, 3.05, 0.28, "综合黑盒：功能 / 安全 / 异常", 11.4, MUTED)
    textbox(slide, 1.0, 2.84, 3.05, 0.28, "压力测试：网关、文章、推荐、当前用户", 11.4, MUTED)
    textbox(slide, 1.0, 3.20, 3.05, 0.28, "推荐专项：构造不同画像用户", 11.4, MUTED)
    textbox(slide, 1.0, 3.56, 3.05, 0.28, "修复复验：针对已发现缺陷", 11.4, MUTED)

    textbox(slide, 4.55, 1.82, 2.05, 0.3, "执行发现", 16, ACCENT, True)
    hline(slide, 4.55, 2.23, 2.7, ACCENT, 1.2)
    textbox(slide, 4.55, 2.48, 3.15, 0.28, "84 项黑盒用例，通过 82 项", 11.4, MUTED)
    textbox(slide, 4.55, 2.84, 3.15, 0.28, "发现注册角色提权风险", 11.4, MUTED)
    textbox(slide, 4.55, 3.20, 3.15, 0.28, "发现 WebSocket 网关转发问题", 11.4, MUTED)
    textbox(slide, 4.55, 3.56, 3.15, 0.28, "发现推荐并发写入冲突", 11.4, MUTED)

    textbox(slide, 8.2, 1.82, 2.05, 0.3, "测试结果", 16, ACCENT, True)
    hline(slide, 8.2, 2.23, 2.7, ACCENT, 1.2)
    textbox(slide, 8.2, 2.48, 3.25, 0.28, "定向复验 7 / 7 通过", 11.4, MUTED)
    textbox(slide, 8.2, 2.84, 3.25, 0.28, "推荐专项 13 / 13 通过", 11.4, MUTED)
    textbox(slide, 8.2, 3.20, 3.25, 0.28, "推荐并发复验 160 次 / 16 并发全 200", 11.4, MUTED)
    textbox(slide, 8.2, 3.56, 3.25, 0.28, "结论：原型闭环可验证，不夸大为生产容量", 11.4, MUTED)

    hline(slide, 1.0, 4.34, 10.6)
    metric(slide, 1.1, 4.75, "84 / 82", "综合黑盒")
    metric(slide, 3.65, 4.75, "7 / 7", "修复复验")
    metric(slide, 6.2, 4.75, "13 / 13", "推荐专项")
    metric(slide, 8.75, 4.75, "160 / 16", "推荐并发")
    soft_band(slide, 1.0, 5.85, 10.65, 0.42, "讲述重点：测试形成了“方案设计 -> 执行发现 -> 修复复验 -> 结论边界”的闭环。", 13.6)


def main():
    prs = Presentation(SRC)
    build_merged_recommendation(prs.slides[5])
    delete_slide(prs, 6)
    # After deleting the old recommendation page, the original testing page moves from slide 9 to slide 8.
    build_testing_flow(prs.slides[7])
    prs.save(OUT)
    print(str(OUT))


if __name__ == "__main__":
    main()
