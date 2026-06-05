# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SRC = Path(r"D:\$RRBBDQK\tmp_ppt_work\defense_ppt_source.pptx")
OUT = Path(r"D:\$RRBBDQK\tmp_ppt_work\defense_ppt_rebuilt.pptx")

FONT = "Microsoft YaHei"
TITLE = RGBColor(0x24, 0x38, 0x46)
ACCENT = RGBColor(0x15, 0x71, 0x6E)
MUTED = RGBColor(0x60, 0x6E, 0x76)
CARD_FILL = RGBColor(0xF8, 0xFC, 0xFB)
ACCENT_FILL = RGBColor(0xE4, 0xF7, 0xF3)
BORDER = RGBColor(0xB3, 0xD9, 0xD3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def remove_shape(shape):
    element = shape._element
    element.getparent().remove(element)


def set_title(slide, title):
    # The original deck uses shape #3 as the slide title on these pages.
    if len(slide.shapes) >= 3 and getattr(slide.shapes[2], "has_text_frame", False):
        tf = slide.shapes[2].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT
        p.font.size = Pt(21)
        p.font.bold = True
        p.font.color.rgb = TITLE


def clear_body(slide):
    # Keep the left bar, footer, slide title and horizontal rule.
    for shape in list(slide.shapes)[4:]:
        remove_shape(shape)


def style_shape(shape, fill=CARD_FILL, line=BORDER, radius=True):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    if radius and hasattr(shape, "adjustments") and len(shape.adjustments) > 0:
        try:
            shape.adjustments[0] = 0.08
        except Exception:
            pass


def add_textbox(slide, x, y, w, h, text, size=18, color=TITLE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_band(slide, x, y, w, h, text, size=20):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    style_shape(shape, ACCENT_FILL, BORDER)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    return shape


def add_card(slide, x, y, w, h, title, body=None, fill=CARD_FILL, title_size=18, body_size=12.5, center=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    style_shape(shape, fill, BORDER)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.06)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if center else MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
    p.font.name = FONT
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    lines = body if isinstance(body, list) else ([body] if body else [])
    for line in lines:
        para = tf.add_paragraph()
        para.text = line
        para.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
        para.space_before = Pt(4)
        para.font.name = FONT
        para.font.size = Pt(body_size)
        para.font.bold = False
        para.font.color.rgb = MUTED
    return shape


def add_metric(slide, x, y, w, h, value, label, note=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    style_shape(shape, ACCENT_FILL, BORDER)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = value
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p2 = tf.add_paragraph()
    p2.text = label
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = FONT
    p2.font.size = Pt(12.5)
    p2.font.color.rgb = MUTED
    if note:
        p3 = tf.add_paragraph()
        p3.text = note
        p3.alignment = PP_ALIGN.CENTER
        p3.font.name = FONT
        p3.font.size = Pt(10.5)
        p3.font.color.rgb = MUTED
    return shape


def add_arrow(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = BORDER
    line.line.width = Pt(1.4)
    return line


def replace_text(slide, replacements):
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text
                for old, new in replacements.items():
                    text = text.replace(old, new)
                run.text = text


def reset_small_stat_text(shape, lines):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT
        p.font.size = Pt(11.5)
        p.font.bold = False
        p.font.color.rgb = MUTED


def fix_front_matter(prs):
    slide2 = prs.slides[1]
    replace_text(
        slide2,
        {
            "用户获取信息渠道中": "用户日常获取信息渠道中",
            "4200万/日": "高频化",
            "网站系统日均面临": "网站系统持续面临",
            "大量自动化威胁": "自动化攻击与风险扫描",
            "35%": "AIGC",
            "新发布网站使用": "生成式 AI 已降低",
            "的比例大幅上述": "生产门槛",
            "74.2%": "内容治理",
            "生成内容在网页内容的占比重剧增": "虚假、低质与垃圾内容增加",
            "通过审计日志保留关键请求与异常线索": "通过 requestId 与审计日志保留关键请求与异常线索",
            "已降低AI生产门槛": "已降低内容生产门槛",
            "AI虚假": "内容风险：虚假",
        },
    )
    ai_count = 0
    for shape in slide2.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if shape.text_frame.text.strip() == "AI":
            ai_count += 1
            shape.text_frame.text = "内容" if ai_count == 1 else "风险"
        compact = shape.text.replace("\n", "").replace("\r", "")
        if "生成式 AI 已降低" in compact and "生产门槛" in compact:
            reset_small_stat_text(shape, ["生成式 AI 已降低", "内容生产门槛"])
        if "虚假、低质" in compact:
            reset_small_stat_text(shape, ["内容风险", "虚假、低质与垃圾内容增加"])

    slide3 = prs.slides[2]
    replace_text(
        slide3,
        {
            "路由转发 / 鉴权": "路由转发 / 限流",
            "限流 / requestId": "requestId 透传",
            "审计入口": "审计查询入口",
        },
    )


def build_slide_4(slide):
    set_title(slide, "设计主线：从发文到追踪")
    clear_body(slide)
    add_band(slide, 0.78, 1.08, 11.65, 0.62, "不是功能堆砌，而是围绕内容治理形成三条闭环", 19)

    cards = [
        ("01  统一入口", ["网关转发", "限流与 requestId"]),
        ("02  内容入库", ["文章先待审核", "互动写业务表"]),
        ("03  AI 治理", ["四项评分", "标签 / 风险 / 决策"]),
        ("04  展示控制", ["先判断能不能看", "再决定排第几"]),
        ("05  推荐排序", ["画像 + 质量 + 新鲜度", "风险与已看惩罚"]),
        ("06  回写追踪", ["推荐事件", "审计日志与复现"]),
    ]
    xs = [0.78, 4.72, 8.66]
    ys = [2.0, 4.08]
    idx = 0
    for y in ys:
        for x in xs:
            title, body = cards[idx]
            add_card(slide, x, y, 3.38, 1.38, title, body, title_size=17.5, body_size=12.5)
            idx += 1

    add_band(slide, 1.05, 5.92, 11.1, 0.5, "讲述抓手：先管住内容，再分发内容，最后留下证据。", 16)


def build_slide_5(slide):
    set_title(slide, "服务拆分：职责边界")
    clear_body(slide)
    add_band(slide, 0.78, 1.08, 11.65, 0.62, "拆分原则：入口统一，业务内聚，AI 与聊天独立演进", 19)
    cards = [
        ("mofukaze", ["Nuxt3 前端", "页面展示与交互"]),
        ("api_gateway", ["统一 API 入口", "转发、限流、requestId"]),
        ("user_service", ["认证与用户资料", "登录、会话、角色"]),
        ("blog_service", ["核心业务服务", "内容、互动、推荐、通知"]),
        ("ai_service", ["AI 分析服务", "评分、标签、风险、Provider"]),
        ("chat_service", ["私信服务", "线程、消息、WebSocket"]),
    ]
    xs = [0.78, 4.72, 8.66]
    ys = [1.98, 3.78]
    idx = 0
    for y in ys:
        for x in xs:
            title, body = cards[idx]
            add_card(slide, x, y, 3.38, 1.38, title, body, title_size=17.2, body_size=12.2)
            idx += 1
    add_metric(slide, 1.05, 5.65, 3.2, 0.72, "6 个服务", "独立镜像与容器")
    add_metric(slide, 5.05, 5.65, 3.2, 0.72, "约 21 张表", "支撑业务闭环")
    add_metric(slide, 9.05, 5.65, 3.2, 0.72, "1 个入口", "网关统一访问")


def build_slide_6(slide):
    set_title(slide, "AI 治理：调用不是终点")
    clear_body(slide)
    add_band(slide, 0.78, 1.08, 11.65, 0.62, "文章先进入待审核；AI 输出必须结构化落库，才能影响后续业务", 18)

    steps = [
        ("触发", ["发文后调用", "blog_service -> ai_service"]),
        ("组织请求", ["messages 风格", "system + user + schema"]),
        ("Provider", ["Mock 稳定测试", "真实模型轮询调用"]),
        ("结构化结果", ["评分 / 标签 / 风险", "决策与建议"]),
    ]
    xs = [0.82, 3.72, 6.62, 9.52]
    for i, (title, body) in enumerate(steps):
        add_card(slide, xs[i], 2.05, 2.38, 1.65, title, body, fill=ACCENT_FILL if i == 2 else CARD_FILL, title_size=18, body_size=11.4, center=True)
        if i < len(steps) - 1:
            add_arrow(slide, xs[i] + 2.43, 2.88, xs[i + 1] - 0.08, 2.88)

    add_card(slide, 0.95, 4.18, 3.45, 1.32, "落库位置", ["article_ai_analysis", "article_ai_tag_on_article", "article.status"], title_size=17, body_size=11.2)
    add_card(slide, 4.85, 4.18, 3.45, 1.32, "业务影响", ["展示控制", "推荐排序", "作者评分与通知"], title_size=17, body_size=11.2)
    add_card(slide, 8.75, 4.18, 3.2, 1.32, "失败策略", ["模型异常不放行", "进入 REVIEW_REQUIRED"], fill=ACCENT_FILL, title_size=17, body_size=11.2)
    add_textbox(slide, 1.0, 5.92, 11.0, 0.42, "这页要讲清楚：AI 不是聊天窗口，而是内容治理链路中的结构化分析服务。", 14.5, MUTED, False, PP_ALIGN.CENTER)


def build_slide_7(slide):
    set_title(slide, "推荐闭环：先过滤，再排序")
    clear_body(slide)
    add_band(slide, 0.78, 1.08, 11.65, 0.62, "推荐不绕过审核：只在可展示候选内做个性化排序", 19)

    add_card(
        slide,
        0.82,
        1.95,
        5.45,
        2.42,
        "可解释推荐公式",
        [
            "tagMatch * 4 + authorAffinity * 2",
            "+ authorQuality * 1.2 + contentQuality + freshness",
            "- riskPenalty - seenPenalty",
        ],
        fill=ACCENT_FILL,
        title_size=19,
        body_size=13,
        center=True,
    )
    add_card(slide, 6.75, 1.95, 2.45, 1.16, "候选过滤", ["PUBLISHED / LOW_PRIORITY", "legalityScore >= 40"], title_size=16.5, body_size=10.8)
    add_card(slide, 9.55, 1.95, 2.45, 1.16, "输入信号", ["画像、AI 标签", "14 天统计、已看记录"], title_size=16.5, body_size=10.8)
    add_card(slide, 6.75, 3.22, 2.45, 1.16, "排序输出", ["推荐列表", "requestId 返回"], title_size=16.5, body_size=10.8)
    add_card(slide, 9.55, 3.22, 2.45, 1.16, "请求记录", ["reco_request_log", "曝光与 seen 回写"], title_size=16.5, body_size=10.8)

    add_band(slide, 0.95, 5.0, 11.0, 0.62, "关键边界：推荐流读取已落库 AI 结果，不会每次推荐都实时请求大模型。", 16)
    add_textbox(slide, 1.0, 5.86, 11.0, 0.42, "讲述顺序：先说候选为什么安全，再说分数为什么可解释，最后说行为如何回写。", 14.2, MUTED, False, PP_ALIGN.CENTER)


def build_slide_8(slide):
    set_title(slide, "画像反馈：行为改变下一次推荐")
    clear_body(slide)
    add_band(slide, 0.78, 1.08, 11.65, 0.62, "用户画像不是手工填写，而是由推荐事件持续计算出来", 19)

    steps = [
        ("展示", ["推荐列表返回", "自动写曝光"]),
        ("行为", ["点击、停留、阅读完成", "点赞、评论、隐藏、举报"]),
        ("事件", ["reco_event", "场景、位置、requestId"]),
        ("画像", ["90 天 / 500 事件", "时间衰减"]),
        ("下次推荐", ["tagVector", "authorAffinity"]),
    ]
    xs = [0.72, 3.05, 5.38, 7.71, 10.04]
    for i, (title, body) in enumerate(steps):
        add_card(slide, xs[i], 2.0, 1.9, 1.65, title, body, fill=ACCENT_FILL if i in (2, 3) else CARD_FILL, title_size=17, body_size=10.1, center=True)
        if i < len(steps) - 1:
            add_arrow(slide, xs[i] + 1.94, 2.83, xs[i + 1] - 0.07, 2.83)

    add_card(slide, 1.0, 4.32, 3.35, 1.25, "正向信号", ["点击、阅读完成、点赞、评论", "提高相关标签和作者权重"], title_size=17, body_size=11.3)
    add_card(slide, 4.95, 4.32, 3.35, 1.25, "负向信号", ["隐藏、举报", "降低相似内容和作者权重"], title_size=17, body_size=11.3)
    add_card(slide, 8.9, 4.32, 3.05, 1.25, "稳定性", ["旧行为逐渐衰减", "避免兴趣长期锁死"], title_size=17, body_size=11.3)
    add_textbox(slide, 1.0, 5.9, 11.0, 0.4, "这页的核心话术：用户不是只接收推荐，用户行为会成为下一次推荐的输入。", 14.2, MUTED, False, PP_ALIGN.CENTER)


def build_slide_9(slide):
    set_title(slide, "安全审计：从请求到证据")
    clear_body(slide)
    add_band(slide, 0.78, 1.08, 11.65, 0.62, "审计不是单独记日志，而是 requestId + 业务动作 + 管理员查询形成证据链", 18)

    cards = [
        ("身份边界", ["access token", "refresh token 仅存哈希"]),
        ("入口边界", ["网关限流", "requestId 透传"]),
        ("权限边界", ["角色判断", "资源归属校验"]),
        ("证据边界", ["audit_log", "用户 / 路径 / 动作 / 结果"]),
    ]
    xs = [0.82, 3.78, 6.74, 9.7]
    for i, (title, body) in enumerate(cards):
        add_card(slide, xs[i], 2.02, 2.45, 1.75, title, body, fill=ACCENT_FILL if i == 3 else CARD_FILL, title_size=17.2, body_size=11.1, center=True)

    add_card(slide, 1.0, 4.35, 3.25, 1.16, "记录什么", ["登录、发文、删除、AI 分析、评论、关注、聊天、异常请求"], title_size=16.4, body_size=10.8)
    add_card(slide, 4.85, 4.35, 3.25, 1.16, "如何查看", ["管理员运维页", "审计列表与摘要统计"], title_size=16.4, body_size=10.8)
    add_card(slide, 8.7, 4.35, 3.25, 1.16, "如何闭环", ["请求可定位", "异常可追踪，结果可复核"], title_size=16.4, body_size=10.8)
    add_textbox(slide, 1.0, 5.9, 11.0, 0.42, "答辩时避免说成“日志很多”，要说成“关键行为可以按请求链路追踪”。", 14.2, MUTED, False, PP_ALIGN.CENTER)


def build_slide_10(slide):
    set_title(slide, "测试证据：闭环是否跑通")
    clear_body(slide)
    add_band(slide, 0.78, 1.08, 11.65, 0.62, "测试目标：验证主流程可运行、问题可发现、修复可复验", 19)

    add_metric(slide, 0.82, 2.02, 2.55, 1.2, "84 / 82", "综合黑盒测试", "通过率 97.62%")
    add_metric(slide, 3.72, 2.02, 2.55, 1.2, "7 / 7", "定向修复复验", "缺陷闭环")
    add_metric(slide, 6.62, 2.02, 2.55, 1.2, "13 / 13", "推荐专项验证", "画像差异有效")
    add_metric(slide, 9.52, 2.02, 2.55, 1.2, "160 / 16", "推荐并发复验", "全 200")

    add_card(slide, 0.95, 3.88, 3.45, 1.45, "发现并修复", ["注册角色提权风险", "WebSocket 网关转发", "推荐并发写入冲突"], title_size=17, body_size=11.3)
    add_card(slide, 4.85, 3.88, 3.45, 1.45, "专项推荐数据", ["3 类画像用户", "36 作者 / 291 文章", "1021 条推荐事件"], title_size=17, body_size=11.3)
    add_card(slide, 8.75, 3.88, 3.2, 1.45, "结论边界", ["证明原型闭环可验证", "不夸大为生产容量测试"], fill=ACCENT_FILL, title_size=17, body_size=11.3)
    add_textbox(slide, 1.0, 5.86, 11.0, 0.42, "系统测试使用 Mock Provider 保证可重复；真实模型接入能力另做补充验证。", 14.2, MUTED, False, PP_ALIGN.CENTER)


def build_slide_11(slide):
    set_title(slide, "总结与致谢")
    clear_body(slide)
    add_band(slide, 0.78, 1.08, 11.65, 0.64, "完成的是一个可运行、可追踪、可解释的社交平台原型", 19)

    cards = [
        ("微服务架构", ["前端、网关、用户", "内容、AI、聊天分离"]),
        ("内容治理", ["AI 评分、标签", "风险与展示状态"]),
        ("推荐闭环", ["画像、事件", "排序与回写"]),
        ("审计闭环", ["requestId", "audit_log 与运维页"]),
    ]
    xs = [0.82, 3.78, 6.74, 9.7]
    for i, (title, body) in enumerate(cards):
        add_card(slide, xs[i], 2.02, 2.45, 1.7, title, body, fill=ACCENT_FILL if i in (1, 2) else CARD_FILL, title_size=17.2, body_size=11.1, center=True)

    add_card(slide, 1.0, 4.22, 10.95, 0.86, "后续优化方向", ["性能优化、单元测试、人工复核、长期推荐效果评估"], fill=CARD_FILL, title_size=16.6, body_size=12, center=True)
    add_textbox(slide, 1.4, 5.55, 10.4, 0.74, "谢谢各位老师", 28, ACCENT, True, PP_ALIGN.CENTER)


def main():
    prs = Presentation(SRC)
    fix_front_matter(prs)
    builders = {
        4: build_slide_4,
        5: build_slide_5,
        6: build_slide_6,
        7: build_slide_7,
        8: build_slide_8,
        9: build_slide_9,
        10: build_slide_10,
        11: build_slide_11,
    }
    for slide_no, builder in builders.items():
        builder(prs.slides[slide_no - 1])
    prs.save(OUT)
    print(str(OUT))


if __name__ == "__main__":
    main()
